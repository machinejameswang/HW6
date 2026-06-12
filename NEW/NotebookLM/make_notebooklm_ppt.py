from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Inches, Pt


NOTEBOOK_DIR = Path(__file__).resolve().parent
BASE_DIR = NOTEBOOK_DIR.parent
PPT_DIR = BASE_DIR / "PPT"
OUTPUT_PPTX = NOTEBOOK_DIR / "hw6_notebooklm_20_slide_source_deck.pptx"
SOURCE_MD = NOTEBOOK_DIR / "hw6_notebooklm_source.md"
FAQ_MD = NOTEBOOK_DIR / "hw6_notebooklm_faq.md"
STUDY_GUIDE_MD = NOTEBOOK_DIR / "hw6_notebooklm_study_guide.md"


COLORS = {
    "navy": RGBColor(31, 41, 55),
    "blue": RGBColor(37, 99, 235),
    "slate": RGBColor(71, 85, 105),
    "gray": RGBColor(241, 245, 249),
    "white": RGBColor(255, 255, 255),
    "green": RGBColor(22, 163, 74),
    "amber": RGBColor(254, 243, 199),
}


SLIDES = [
    {
        "title": "HW6 Project Overview",
        "key_facts": [
            "Project: Kaggle 50 Startups Profit Prediction",
            "Methodology: CRISP-DM",
            "Library: Scikit-learn",
            "Task type: supervised regression",
        ],
        "why": "This slide introduces the project scope and frames the work as a complete machine learning workflow.",
        "narration": "This project predicts startup profit using the Kaggle 50 Startups dataset. It follows CRISP-DM and uses Scikit-learn regression models.",
    },
    {
        "title": "Business Objective",
        "key_facts": [
            "Goal: predict Profit from spending variables",
            "Use case: resource allocation and scenario planning",
            "Audience: founders, investors, analysts, and managers",
            "Interpretation: predictive association, not causality",
        ],
        "why": "The model supports better business discussion about limited startup resources.",
        "narration": "The business objective is to estimate profit from startup spending patterns. The result supports decision-making but should not be interpreted as causal proof.",
    },
    {
        "title": "Dataset Structure",
        "key_facts": [
            "File: 50_Startups_dataset.csv",
            "Rows: 50",
            "Target: Profit",
            "Features: R&D Spend, Administration, Marketing Spend, State",
        ],
        "why": "The dataset is small, so cross-validation is important.",
        "narration": "The dataset has 50 records and five useful columns. Profit is the target variable, and the remaining columns are predictors.",
    },
    {
        "title": "CRISP-DM Six Steps",
        "key_facts": [
            "1. Business Understanding",
            "2. Data Understanding",
            "3. Data Preparation",
            "4. Modeling",
            "5. Evaluation",
            "6. Deployment",
        ],
        "why": "CRISP-DM connects business questions, data analysis, modeling, evaluation, and deployment.",
        "narration": "The workflow follows the six steps of CRISP-DM, giving the project a clear structure from problem definition to deployment simulation.",
    },
    {
        "title": "Data Quality",
        "key_facts": [
            "Dataset shape after cleaning: 50 rows, 5 columns",
            "Missing values: 0",
            "Duplicate rows: 0",
            "States: New York, California, Florida",
        ],
        "why": "Clean input data makes the modeling workflow simpler, but the small sample size remains a limitation.",
        "narration": "The dataset has no missing values and no duplicate rows. The state categories are New York, California, and Florida.",
    },
    {
        "title": "Feature Interpretation",
        "key_facts": [
            "R&D Spend: product innovation and technical capability",
            "Marketing Spend: market expansion and customer acquisition",
            "Administration: operating scale and management structure",
            "State: regional auxiliary factor",
        ],
        "why": "Feature meaning helps explain why a model is useful for business decisions.",
        "narration": "Each feature has a business meaning. R&D is linked to innovation, Marketing to market reach, Administration to operating scale, and State to regional context.",
    },
    {
        "title": "Correlation Evidence",
        "key_facts": [
            "R&D Spend correlation with Profit: 0.9729",
            "Marketing Spend correlation with Profit: 0.7478",
            "Administration correlation with Profit: 0.2007",
            "R&D Spend has the strongest direct linear association",
        ],
        "why": "Correlation gives an initial signal about which numeric features are most associated with Profit.",
        "narration": "R&D Spend has the strongest correlation with Profit. Marketing is also positive but weaker, while Administration is much weaker.",
    },
    {
        "title": "Data Preparation",
        "key_facts": [
            "Separate X features and y target",
            "Use Profit as target",
            "Encode State using OneHotEncoder",
            "Use ColumnTransformer and Pipeline",
            "Use train_test_split with test_size=0.2 and random_state=42",
        ],
        "why": "A Scikit-learn Pipeline keeps preprocessing and modeling reproducible.",
        "narration": "The preparation step separates features and target, encodes the categorical State variable, and uses a Pipeline for reproducibility.",
    },
    {
        "title": "Modeling Strategy",
        "key_facts": [
            "Primary model: LinearRegression",
            "Reason: target is continuous",
            "Reason: dataset is small",
            "Reason: model is interpretable",
        ],
        "why": "Linear Regression is simple, explainable, and appropriate for a teaching-focused regression project.",
        "narration": "Linear Regression was selected because it is interpretable and suitable for a small regression dataset.",
    },
    {
        "title": "Model Experiments",
        "key_facts": [
            "Model 1: R&D only",
            "Model 2: R&D + Marketing",
            "Model 3: numerical features",
            "Model 4: all features",
        ],
        "why": "Comparing feature sets tests whether more variables actually improve predictive performance.",
        "narration": "The project compares four models to test the value of R&D, Marketing, Administration, and State.",
    },
    {
        "title": "Evaluation Metrics",
        "key_facts": [
            "Train-test metrics: R2, MAE, RMSE",
            "Cross-validation: 5-fold KFold",
            "CV metrics: CV R2 Mean, CV R2 Std, CV RMSE Mean, CV RMSE Std",
            "Small dataset requires cross-validation",
        ],
        "why": "Cross-validation reduces reliance on a single train-test split.",
        "narration": "The models are evaluated using both train-test metrics and 5-fold cross-validation because the dataset has only 50 rows.",
    },
    {
        "title": "Model Comparison Results",
        "key_facts": [
            "R&D only: Test R2 0.9265, CV R2 Mean 0.9374",
            "R&D + Marketing: Test R2 0.9168, CV R2 Mean 0.9389",
            "Numerical Features: Test R2 0.9001, CV R2 Mean 0.9338",
            "All Features: Test R2 0.8987, CV R2 Mean 0.9279",
        ],
        "why": "R&D + Marketing has the best CV R2 Mean, while the full model performs worse.",
        "narration": "The model comparison shows that R&D plus Marketing has the highest cross-validation R2 mean, while the all-feature model performs worse.",
    },
    {
        "title": "Feature Selection Methods",
        "key_facts": [
            "Pearson correlation",
            "SelectKBest F-regression",
            "Recursive Feature Elimination",
            "LASSO",
            "Random Forest importance",
        ],
        "why": "Using several methods makes the feature conclusion more robust.",
        "narration": "Several feature selection methods were used to avoid relying on only one technique.",
    },
    {
        "title": "Feature Selection Consensus",
        "key_facts": [
            "R&D Spend is dominant",
            "Marketing Spend is secondary",
            "Administration has limited incremental value",
            "State has very limited value in this small dataset",
        ],
        "why": "The same conclusion appears across statistical, model-based, regularized, and tree-based methods.",
        "narration": "The feature selection methods consistently identify R&D Spend as the strongest feature, with Marketing as a secondary feature.",
    },
    {
        "title": "Final Model Recommendation",
        "key_facts": [
            "Main model: Profit = f(R&D Spend, Marketing Spend)",
            "Benchmark model: Profit = f(R&D Spend)",
            "Do not choose the most complex model automatically",
            "Prefer performance, stability, simplicity, and interpretability",
        ],
        "why": "The final model balances predictive performance with business meaning.",
        "narration": "The recommended main model uses R&D Spend and Marketing Spend. R&D alone remains the benchmark.",
    },
    {
        "title": "Deployment Simulation",
        "key_facts": [
            "Sample R&D Spend: 120,000",
            "Sample Administration: 130,000",
            "Sample Marketing Spend: 250,000",
            "Sample State: New York",
            "Predicted Profit: 150,042.94",
        ],
        "why": "The deployment step demonstrates how the final pipeline can make a prediction for a new startup.",
        "narration": "The deployment simulation predicts profit for a sample startup and saves the model pipeline for reuse.",
    },
    {
        "title": "Expert Consensus",
        "key_facts": [
            "Marketing expert: Marketing supports market reach",
            "Top CEO: model must be explainable in business meetings",
            "Policy expert: State needs more data before strong claims",
            "R&D expert: R&D is the core innovation signal",
        ],
        "why": "The expert panel makes the model interpretation more realistic and business-facing.",
        "narration": "The expert panel agrees that R&D is the core signal, Marketing adds business context, and State should not be overinterpreted.",
    },
    {
        "title": "Limitations",
        "key_facts": [
            "Only 50 observations",
            "Observational data",
            "Possible multicollinearity between R&D and Marketing",
            "State has limited samples per category",
            "Results are predictive associations, not causal conclusions",
        ],
        "why": "Limitations prevent overclaiming and guide future improvements.",
        "narration": "The dataset is small and observational, so the results should be interpreted carefully.",
    },
    {
        "title": "Artifacts Produced",
        "key_facts": [
            "V1 and V2 CRISP-DM Python scripts",
            "Technical charts and CSV summaries",
            "Final model pipeline",
            "20-slide PPT with notes and male narration",
            "HeyGen and Hyperframes packages",
        ],
        "why": "The project includes code, reports, visuals, models, and presentation materials.",
        "narration": "The final project includes scripts, charts, summaries, models, and presentation files.",
    },
    {
        "title": "Final Conclusion",
        "key_facts": [
            "R&D Spend is the strongest predictor",
            "Marketing Spend adds secondary business value",
            "Administration and State add limited value here",
            "Use R&D + Marketing as the main model",
            "Keep R&D only as the benchmark",
        ],
        "why": "This conclusion is supported by model evaluation and feature selection evidence.",
        "narration": "The final conclusion is to use R&D Spend plus Marketing Spend as the main model, while keeping R&D Spend alone as the benchmark.",
    },
]


def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color or COLORS["navy"]
    return box


def add_bullet_section(slide, title, bullets, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["gray"]
    shape.line.color.rgb = COLORS["blue"]
    shape.text_frame.clear()
    shape.text_frame.margin_left = Inches(0.15)
    shape.text_frame.margin_right = Inches(0.15)
    shape.text_frame.margin_top = Inches(0.1)
    p = shape.text_frame.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(17)
    p.font.color.rgb = COLORS["blue"]
    for bullet in bullets:
        bp = shape.text_frame.add_paragraph()
        bp.text = bullet
        bp.level = 0
        bp.font.size = Pt(14)
        bp.font.color.rgb = COLORS["navy"]


def add_why_box(slide, text):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.75),
        Inches(5.85),
        Inches(11.8),
        Inches(0.85),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["amber"]
    shape.line.color.rgb = COLORS["green"]
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "Why it matters: " + text
    p.font.size = Pt(16)
    p.font.color.rgb = COLORS["navy"]
    p.font.bold = True


def add_notes(slide, data):
    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    notes.text = (
        f"NotebookLM source notes\n\n"
        f"Slide title: {data['title']}\n\n"
        f"Key facts:\n- " + "\n- ".join(data["key_facts"]) + "\n\n"
        f"Why it matters: {data['why']}\n\n"
        f"Narration: {data['narration']}"
    )


def build_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for idx, data in enumerate(SLIDES, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = COLORS["white"]
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.18))
        bar.fill.solid()
        bar.fill.fore_color.rgb = COLORS["blue"]
        bar.line.fill.background()
        add_textbox(slide, Inches(0.75), Inches(0.45), Inches(11.8), Inches(0.45), f"{idx}. {data['title']}", 28, True)
        add_bullet_section(slide, "Key facts", data["key_facts"], Inches(0.8), Inches(1.25), Inches(5.9), Inches(4.25))
        add_bullet_section(slide, "Narration summary", [data["narration"]], Inches(7.0), Inches(1.25), Inches(5.55), Inches(4.25))
        add_why_box(slide, data["why"])
        add_notes(slide, data)
    prs.save(OUTPUT_PPTX)


def write_markdown_sources():
    lines = [
        "# HW6 NotebookLM Source",
        "",
        "This source summarizes the HW6 Kaggle 50 Startups CRISP-DM project for NotebookLM.",
        "",
        "## Core Conclusion",
        "",
        "Use `R&D Spend + Marketing Spend` as the main explainable prediction model.",
        "Keep `R&D Spend` alone as the benchmark model.",
        "Interpret the results as predictive associations, not causal conclusions.",
        "",
    ]
    for idx, data in enumerate(SLIDES, start=1):
        lines.extend([
            f"## Slide {idx}: {data['title']}",
            "",
            "### Key facts",
            "",
        ])
        lines.extend([f"- {fact}" for fact in data["key_facts"]])
        lines.extend([
            "",
            "### Why it matters",
            "",
            data["why"],
            "",
            "### Narration",
            "",
            data["narration"],
            "",
        ])
    SOURCE_MD.write_text("\n".join(lines), encoding="utf-8")

    FAQ_MD.write_text(
        """# HW6 NotebookLM FAQ

## What is the goal of the project?

The goal is to predict startup Profit using R&D Spend, Administration, Marketing Spend, and State.

## What type of machine learning problem is this?

It is a supervised regression problem because Profit is a continuous numeric target.

## Why is R&D Spend important?

R&D Spend is the strongest predictor across correlation, cross-validation, and feature selection methods.

## Why include Marketing Spend?

Marketing Spend adds secondary business value and gives the model a stronger market-expansion interpretation.

## Why not use all features?

The all-feature model has lower CV R2 Mean than the R&D + Marketing model, so more features do not improve performance here.

## Can the model prove causality?

No. The dataset is observational and has only 50 records, so the results are predictive associations, not causal conclusions.

## What is the final model?

The final main model is `Profit = f(R&D Spend, Marketing Spend)`.

## What is the benchmark model?

The benchmark model is `Profit = f(R&D Spend)`.
""",
        encoding="utf-8",
    )

    STUDY_GUIDE_MD.write_text(
        """# HW6 NotebookLM Study Guide

## One-sentence summary

This CRISP-DM project predicts startup Profit and concludes that R&D Spend is the dominant predictor, with Marketing Spend as a useful secondary feature.

## Must-know concepts

- CRISP-DM
- Supervised regression
- Linear Regression
- Train-test split
- 5-fold cross-validation
- OneHotEncoder
- ColumnTransformer
- Pipeline
- R2, MAE, RMSE
- Feature selection

## Final model

```text
Profit = f(R&D Spend, Marketing Spend)
```

## Benchmark model

```text
Profit = f(R&D Spend)
```

## Key numbers

- R&D only CV R2 Mean: 0.9374
- R&D + Marketing CV R2 Mean: 0.9389
- All features CV R2 Mean: 0.9279
- Sample predicted Profit: 150,042.94

## Main caution

The dataset has only 50 rows, so use cross-validation and avoid causal claims.
""",
        encoding="utf-8",
    )


def main():
    build_ppt()
    write_markdown_sources()
    print(f"Saved NotebookLM PPT: {OUTPUT_PPTX}")
    print(f"Saved NotebookLM source: {SOURCE_MD}")
    print(f"Saved NotebookLM FAQ: {FAQ_MD}")
    print(f"Saved NotebookLM study guide: {STUDY_GUIDE_MD}")


if __name__ == "__main__":
    main()
