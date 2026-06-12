from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


PPT_DIR = Path(__file__).resolve().parent
INPUT_PPTX = PPT_DIR / "hw6_50_startups_crispdm_20_slides_with_notes.pptx"
OUTPUT_PPTX = PPT_DIR / "hw6_50_startups_crispdm_20_slides_male_voice_embedded.pptx"
AUDIO_DIR = PPT_DIR / "audio"


def main():
    prs = Presentation(INPUT_PPTX)

    for index, slide in enumerate(prs.slides, start=1):
        audio_path = AUDIO_DIR / f"slide_{index:02d}_male_narration.wav"
        if not audio_path.exists():
            raise FileNotFoundError(f"Missing audio file: {audio_path}")

        # Small media icon in the bottom-right corner. In PowerPoint, click it
        # to play the slide narration. Speaker notes remain available too.
        slide.shapes.add_movie(
            str(audio_path),
            Inches(12.25),
            Inches(6.25),
            Inches(0.35),
            Inches(0.35),
            mime_type="audio/wav",
        )

    prs.save(OUTPUT_PPTX)
    print(f"Saved embedded-audio PowerPoint: {OUTPUT_PPTX}")


if __name__ == "__main__":
    main()
