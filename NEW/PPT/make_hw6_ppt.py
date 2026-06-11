from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


PPT_DIR = Path(__file__).resolve().parent
BASE_DIR = PPT_DIR.parent
OUTPUT_DIR = BASE_DIR / "outputs"
CHART_DIR = OUTPUT_DIR / "algorithm_charts"
PPTX_PATH = PPT_DIR / "hw6_50_startups_crispdm_20_slides_with_notes.pptx"
NARRATION_PATH = PPT_DIR / "hw6_20_slide_voiceover_script.txt"


COLORS = {
    "navy": RGBColor(31, 41, 55),
    "slate": RGBColor(71, 85, 105),
    "blue": RGBColor(37, 99, 235),
    "sky": RGBColor(224, 242, 254),
    "green": RGBColor(22, 163, 74),
    "mint": RGBColor(220, 252, 231),
    "orange": RGBColor(234, 88, 12),
    "amber": RGBColor(254, 243, 199),
    "purple": RGBColor(124, 58, 237),
    "lavender": RGBColor(237, 233, 254),
    "rose": RGBColor(254, 226, 226),
    "white": RGBColor(255, 255, 255),
    "black": RGBColor(15, 23, 42),
    "gray": RGBColor(241, 245, 249),
}


SLIDES = [
    {
        "title": "HW6: Kaggle 50 Startups Profit Prediction",
        "subtitle": "CRISP-DM + Scikit-learn Regression Project",
        "bullets": ["Prepared by: James Wang", "Target: startup Profit", "Final model: R&D Spend + Marketing Spend"],
        "visual": OUTPUT_DIR / "hw6_crispdm_excalidraw_infographic_a4_2k_upscaled.png",
        "notes": "Welcome to my HW6 project presentation. In this project, I use the Kaggle 50 Startups dataset to predict startup profit. The analysis follows the CRISP-DM process and uses Scikit-learn regression models to compare different feature sets.",
    },
    {
        "title": "Project Objective",
        "bullets": ["Predict startup Profit", "Use business spending variables", "Build an interpretable regression model", "Support business decision-making"],
        "callout": "Business spending -> predicted Profit",
        "notes": "The main objective is to predict startup profit based on spending information. The model is designed to be interpretable, so it can support decisions by founders, investors, analysts, and managers.",
    },
    {
        "title": "Dataset Overview",
        "bullets": ["Dataset: 50_Startups_dataset.csv", "Records: 50", "Target: Profit", "Features: R&D Spend, Administration, Marketing Spend, State", "Problem type: supervised regression"],
        "callout": "50 rows | 5 columns | clean dataset",
        "notes": "The dataset contains 50 startup records. The target variable is Profit. The predictors include R&D spending, administrative spending, marketing spending, and state. Since Profit is a continuous numeric value, this is a supervised regression problem.",
    },
    {
        "title": "CRISP-DM Workflow",
        "bullets": ["1. Business Understanding", "2. Data Understanding", "3. Data Preparation", "4. Modeling", "5. Evaluation", "6. Deployment"],
        "callout": "Plan -> Model -> Evaluate -> Deploy",
        "notes": "This project follows the six steps of CRISP-DM. This structure helps connect the business question, data exploration, model building, evaluation, and final deployment simulation into one complete workflow.",
    },
    {
        "title": "Business Understanding",
        "bullets": ["Business question: which spending patterns are associated with higher profit?", "Startup resources are limited", "Model supports planning and resource allocation", "Important: prediction, not causality"],
        "callout": "Use association language, not causal claims.",
        "notes": "The business question is which spending patterns are associated with higher profit. Because the dataset is observational and small, I avoid claiming causality. The model can identify predictive associations, but it cannot prove that one feature directly causes profit to increase.",
    },
    {
        "title": "Expert Panel View",
        "bullets": ["Marketing expert: market expansion", "Top CEO: decision usefulness", "Governor / policy expert: regional factors", "R&D expert: product innovation"],
        "callout": "Four perspectives, one model decision.",
        "notes": "To make the interpretation more business-oriented, I use four expert perspectives. The marketing expert focuses on customer acquisition, the CEO focuses on decision usefulness, the policy expert considers regional factors, and the R&D expert focuses on product innovation.",
    },
    {
        "title": "Data Understanding",
        "bullets": ["Shape: 50 rows, 5 useful columns", "Missing values: 0", "Duplicate rows: 0", "States: New York, California, Florida"],
        "callout": "Data quality check passed.",
        "notes": "The dataset is clean. It has 50 rows and 5 useful columns after removing the exported index column. There are no missing values and no duplicate rows. The state variable contains New York, California, and Florida.",
    },
    {
        "title": "Feature Meaning",
        "bullets": ["R&D Spend: product innovation and technical capability", "Marketing Spend: brand exposure and customer acquisition", "Administration: operating scale and management structure", "State: regional auxiliary factor"],
        "callout": "Feature meaning guides interpretation.",
        "notes": "Each feature has a business meaning. R&D spending reflects product development and innovation. Marketing spending reflects market expansion and customer acquisition. Administration may represent operating scale. State may represent regional conditions, but the sample size is too small for strong regional conclusions.",
    },
    {
        "title": "Correlation Results",
        "bullets": ["R&D Spend: 0.9729", "Marketing Spend: 0.7478", "Administration: 0.2007", "R&D is the strongest direct linear signal"],
        "visual": CHART_DIR / "01_pearson_correlation.png",
        "notes": "The correlation results show that R&D Spend has the strongest direct linear relationship with Profit. Marketing Spend is also positively associated with Profit, but it is weaker than R&D. Administration has a much weaker correlation.",
    },
    {
        "title": "Data Preparation",
        "bullets": ["Separate X and y", "Use Profit as target", "Encode State with OneHotEncoder", "Use ColumnTransformer and Pipeline", "Train-test split: 80/20 with random_state=42"],
        "callout": "Reproducible sklearn Pipeline",
        "notes": "For data preparation, I separate the features and target variable. The categorical State feature is encoded with OneHotEncoder, not label encoding. I use ColumnTransformer and Pipeline to keep preprocessing and modeling in one reproducible workflow.",
    },
    {
        "title": "Modeling Strategy",
        "bullets": ["Primary model: LinearRegression", "Target is continuous", "Dataset is small", "Model is interpretable", "Suitable for CRISP-DM teaching and reporting"],
        "callout": "Simple model, clear explanation.",
        "notes": "The primary model is Linear Regression. This model is appropriate because the target is continuous, the dataset is small, and the results are easier to explain in a business and educational setting.",
    },
    {
        "title": "Model Experiments",
        "bullets": ["Model 1: R&D only", "Model 2: R&D + Marketing", "Model 3: Numerical features", "Model 4: All features"],
        "callout": "Test whether more features actually help.",
        "notes": "I compare four feature sets. First, I test R&D alone. Second, I add Marketing. Third, I include all numerical features. Finally, I include all features, including State. This helps test whether additional variables actually improve prediction.",
    },
    {
        "title": "Evaluation Metrics",
        "bullets": ["R2 Score", "MAE", "RMSE", "CV R2 Mean and Std", "CV RMSE Mean and Std", "5-fold cross-validation"],
        "callout": "Small dataset -> use cross-validation.",
        "notes": "The models are evaluated with R2, MAE, and RMSE on the train-test split. Since the dataset has only 50 rows, I also use 5-fold cross-validation to check model stability and reduce reliance on one split.",
    },
    {
        "title": "Model Comparison",
        "bullets": ["R&D only: Test R2 0.9265 | CV R2 0.9374", "R&D + Marketing: Test R2 0.9168 | CV R2 0.9389", "Numerical features: Test R2 0.9001 | CV R2 0.9338", "All features: Test R2 0.8987 | CV R2 0.9279"],
        "callout": "Best CV R2 Mean: R&D + Marketing",
        "notes": "The comparison shows that R&D only performs very strongly. However, R&D plus Marketing has the highest cross-validation R2 mean, at 0.9389. The numerical and all-feature models perform worse, which shows that more features do not automatically improve the model.",
    },
    {
        "title": "Feature Selection Evidence",
        "bullets": ["Pearson correlation", "SelectKBest F-regression", "Recursive Feature Elimination", "LASSO regularization", "Random Forest importance"],
        "callout": "Consensus: R&D Spend dominates.",
        "notes": "Multiple feature selection methods support the same conclusion. Pearson correlation, SelectKBest, RFE, LASSO, and Random Forest importance all show that R&D Spend is the dominant predictive feature.",
    },
    {
        "title": "Random Forest Importance",
        "bullets": ["R&D Spend: 0.9228", "Marketing Spend: 0.0670", "Administration: 0.0072", "State total: about 0.0030"],
        "visual": CHART_DIR / "05_random_forest_importance.png",
        "notes": "The Random Forest feature importance result also confirms the same pattern. R&D Spend contributes most of the predictive signal. Marketing is second, while Administration and State contribute very little in this dataset.",
    },
    {
        "title": "Final Model Decision",
        "bullets": ["Main model: Profit = f(R&D Spend, Marketing Spend)", "Benchmark: Profit = f(R&D Spend)", "Reason: best CV R2 Mean plus business meaning", "Avoid selecting the most complex model automatically"],
        "callout": "Final model: R&D + Marketing",
        "notes": "The final main model uses R&D Spend and Marketing Spend. R&D alone is kept as the simplest benchmark. This choice balances predictive performance, business meaning, simplicity, and interpretability.",
    },
    {
        "title": "Deployment Simulation",
        "bullets": ["R&D Spend: 120,000", "Administration: 130,000", "Marketing Spend: 250,000", "State: New York", "Predicted Profit: 150,042.94"],
        "callout": "Learning-project deployment, not full production.",
        "notes": "For deployment simulation, I use a sample startup with R&D Spend of 120,000, Administration of 130,000, Marketing Spend of 250,000, and State as New York. The final model predicts a profit of about 150,042.94.",
    },
    {
        "title": "Business Interpretation",
        "bullets": ["Prioritize product innovation", "Use Marketing to support market reach", "Administration and State add limited value here", "Do not overclaim causality"],
        "callout": "Product innovation is the core signal.",
        "notes": "The business interpretation is that product innovation appears to be the strongest signal associated with profit. Marketing can support market reach, but its incremental value is small. Administration and State were tested, but they do not add meaningful predictive value in this small dataset.",
    },
    {
        "title": "Final Conclusion",
        "bullets": ["Use R&D Spend + Marketing Spend as the main explainable model", "Keep R&D Spend as the benchmark", "Treat results as predictive associations", "Dataset is small: 50 observations"],
        "visual": OUTPUT_DIR / "hw6_crispdm_excalidraw_infographic_a4_2k_upscaled.png",
        "notes": "In conclusion, the recommended model is R&D Spend plus Marketing Spend. R&D Spend remains the strongest predictor, while Marketing provides useful business context. Because the dataset contains only 50 observations, the results should be treated as predictive associations rather than causal conclusions. Thank you.",
    },
]


def add_textbox(slide, left, top, width, height, text, font_size=20, color=None, bold=False, align=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    box.text_frame.clear()
    paragraph = box.text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.size = Pt(font_size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color or COLORS["black"]
    if align is not None:
        paragraph.alignment = align
    return box


def add_bullets(slide, bullets):
    box = slide.shapes.add_textbox(Inches(0.65), Inches(1.55), Inches(5.75), Inches(4.65))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(22 if len(bullet) < 48 else 19)
        p.font.color.rgb = COLORS["navy"]
        p.space_after = Pt(8)


def add_callout(slide, text):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(7.0),
        Inches(2.25),
        Inches(5.35),
        Inches(2.0),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["amber"]
    shape.line.color.rgb = COLORS["orange"]
    shape.line.width = Pt(2)
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = COLORS["black"]
    p.alignment = PP_ALIGN.CENTER


def add_visual(slide, image_path):
    if not image_path or not Path(image_path).exists():
        return
    slide.shapes.add_picture(
        str(image_path),
        Inches(7.05),
        Inches(1.45),
        width=Inches(5.05),
    )


def add_footer(slide, slide_number):
    add_textbox(
        slide,
        Inches(0.65),
        Inches(6.92),
        Inches(7.0),
        Inches(0.25),
        "HW6 | Kaggle 50 Startups | CRISP-DM",
        font_size=9,
        color=COLORS["slate"],
    )
    add_textbox(
        slide,
        Inches(12.0),
        Inches(6.92),
        Inches(0.55),
        Inches(0.25),
        str(slide_number),
        font_size=9,
        color=COLORS["slate"],
        align=PP_ALIGN.RIGHT,
    )


def add_notes(slide, notes):
    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.clear()
    notes_frame.text = notes


def add_slide(prs, idx, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS["white"]

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.18))
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLORS["blue"]
    accent.line.fill.background()

    add_textbox(
        slide,
        Inches(0.65),
        Inches(0.45),
        Inches(11.8),
        Inches(0.55),
        data["title"],
        font_size=30,
        color=COLORS["navy"],
        bold=True,
    )

    if data.get("subtitle"):
        add_textbox(
            slide,
            Inches(0.68),
            Inches(1.02),
            Inches(8.0),
            Inches(0.35),
            data["subtitle"],
            font_size=17,
            color=COLORS["slate"],
        )

    add_bullets(slide, data["bullets"])

    if data.get("visual"):
        add_visual(slide, data["visual"])
    else:
        add_callout(slide, data["callout"])

    note_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.65),
        Inches(6.18),
        Inches(12.0),
        Inches(0.55),
    )
    note_box.fill.solid()
    note_box.fill.fore_color.rgb = COLORS["gray"]
    note_box.line.color.rgb = COLORS["slate"]
    note_box.line.width = Pt(0.5)
    note_tf = note_box.text_frame
    note_tf.clear()
    note_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = note_tf.paragraphs[0]
    p.text = "Speaker notes included in PowerPoint Notes pane."
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS["slate"]
    p.alignment = PP_ALIGN.CENTER

    add_footer(slide, idx)
    add_notes(slide, data["notes"])


def write_narration_file():
    lines = []
    for i, slide in enumerate(SLIDES, start=1):
        lines.append(f"Slide {i}: {slide['title']}")
        lines.append(slide["notes"])
        lines.append("")
    NARRATION_PATH.write_text("\n".join(lines), encoding="utf-8")


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for index, slide_data in enumerate(SLIDES, start=1):
        add_slide(prs, index, slide_data)

    prs.save(PPTX_PATH)
    write_narration_file()
    print(f"Saved PowerPoint: {PPTX_PATH}")
    print(f"Saved narration: {NARRATION_PATH}")


if __name__ == "__main__":
    main()
